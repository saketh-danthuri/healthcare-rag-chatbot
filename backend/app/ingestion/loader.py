"""
loader.py - Multi-Format Document Loader
==========================================
WHY: Your Docs folder has 160+ files in PDF, PPTX, DOC, and TXT formats.
     Each format needs a different extraction library. This module provides
     a single `load_all_documents()` function that handles all formats and
     returns a uniform list of Document objects with rich metadata.

HOW IT WORKS:
  1. Walks through all subfolders (Files, Files_1, Knowledge, Temp)
  2. Detects file type by extension
  3. Uses the best library for each format:
     - PDF: PyMuPDF4LLM (fast, preserves tables as markdown)
     - PPTX: python-pptx (extracts slide text + notes)
     - DOCX: python-docx (preserves paragraphs)
     - DOC: textract fallback (legacy Word format)
     - TXT: built-in Python
  4. Attaches metadata: source file, folder, doc type, page numbers
"""

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)


@dataclass
class Document:
    """A single document with its extracted text and metadata.

    This is our internal representation - not tied to LangChain's Document
    class so we stay framework-agnostic. The chunker will split these into
    smaller pieces for embedding.
    """

    content: str
    metadata: dict = field(default_factory=dict)


def _extract_job_id(filename: str) -> str:
    """Extract the job ID from runbook filenames like 'ATL101Y_1D00_...' or 'CFT303A_1D00_...'.

    WHY: Job IDs are the primary way ops teams search for runbooks.
    Storing this as metadata lets us do precise filtering like
    "show me all chunks from CFT303A".
    """
    match = re.match(r"^([A-Z]{2,4}\d{2,5}[A-Z]?)", filename)
    return match.group(1) if match else ""


def _classify_doc_type(filepath: Path) -> str:
    """Classify document type based on folder and filename patterns.

    WHY: Different doc types need different chunking strategies.
    Runbooks have structured sections (Purpose, Steps, Failure Instructions).
    Training PPTs need slide-by-slide extraction. Knowledge docs are freeform.
    """
    folder = filepath.parent.name.lower()
    filename = filepath.name.lower()

    if "training" in filename or "ppt" in filename:
        return "training"
    if folder in ("files", "files_1", "temp"):
        return "runbook"
    if folder == "knowledge":
        return "knowledge"
    return "general"


def load_pdf(filepath: Path) -> list[Document]:
    """Load a PDF file using PyMuPDF4LLM for high-quality text extraction.

    WHY PyMuPDF4LLM over other libraries:
    - 10x faster than pdfplumber
    - Preserves tables as markdown (critical for runbook tables)
    - Handles complex layouts (multi-column, embedded images)
    - Extracts page-by-page with page numbers
    """
    try:
        import pymupdf4llm

        md_text = pymupdf4llm.to_markdown(str(filepath), page_chunks=True)

        documents = []
        for page_data in md_text:
            text = page_data.get("text", "").strip()
            if not text:
                continue

            page_num = page_data.get("metadata", {}).get("page", 0)

            documents.append(
                Document(
                    content=text,
                    metadata={
                        "source_file": filepath.name,
                        "source_path": str(filepath),
                        "folder": filepath.parent.name,
                        "doc_type": _classify_doc_type(filepath),
                        "job_id": _extract_job_id(filepath.stem),
                        "page_number": page_num,
                        "file_type": "pdf",
                    },
                )
            )

        logger.info(f"Loaded PDF: {filepath.name} ({len(documents)} pages)")
        return documents

    except Exception as e:
        logger.error(f"Failed to load PDF {filepath.name}: {e}")
        # Fallback: try basic PyMuPDF extraction
        try:
            import pymupdf

            doc = pymupdf.open(str(filepath))
            documents = []
            for page_num, page in enumerate(doc):
                text = page.get_text().strip()
                if text:
                    documents.append(
                        Document(
                            content=text,
                            metadata={
                                "source_file": filepath.name,
                                "source_path": str(filepath),
                                "folder": filepath.parent.name,
                                "doc_type": _classify_doc_type(filepath),
                                "job_id": _extract_job_id(filepath.stem),
                                "page_number": page_num + 1,
                                "file_type": "pdf",
                            },
                        )
                    )
            doc.close()
            logger.info(
                f"Loaded PDF (fallback): {filepath.name} ({len(documents)} pages)"
            )
            return documents
        except Exception as e2:
            logger.error(
                f"Fallback PDF extraction also failed for {filepath.name}: {e2}"
            )
            return []


def load_pptx(filepath: Path) -> list[Document]:
    """Load a PowerPoint file, extracting text from each slide.

    WHY slide-by-slide: Each slide in your COPS Training PPTs covers a
    specific topic (e.g., "CLMU Keywords", "Tidal Job Monitoring"). Chunking
    by slide preserves topic boundaries naturally.
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(filepath))
        documents = []

        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []

            # Extract text from all shapes (text boxes, titles, tables)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Extract table content
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text.strip(" |"):
                            texts.append(row_text)

            # Also check slide notes (often contain detailed procedures)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"\nSlide Notes: {notes}")

            slide_content = "\n".join(texts).strip()
            if slide_content:
                documents.append(
                    Document(
                        content=slide_content,
                        metadata={
                            "source_file": filepath.name,
                            "source_path": str(filepath),
                            "folder": filepath.parent.name,
                            "doc_type": _classify_doc_type(filepath),
                            "job_id": _extract_job_id(filepath.stem),
                            "page_number": slide_num,
                            "file_type": "pptx",
                        },
                    )
                )

        logger.info(f"Loaded PPTX: {filepath.name} ({len(documents)} slides)")
        return documents

    except Exception as e:
        logger.error(f"Failed to load PPTX {filepath.name}: {e}")
        return []


def load_docx(filepath: Path) -> list[Document]:
    """Load a Word .docx file, preserving paragraph structure.

    WHY python-docx: It reads the XML structure of .docx files directly,
    preserving headings, bold text, and paragraph breaks. This helps the
    chunker identify section boundaries.
    """
    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(str(filepath))
        texts = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Mark headings for the chunker to use as split points
                if para.style and para.style.name.startswith("Heading"):
                    texts.append(f"\n## {text}")
                else:
                    texts.append(text)

        # Also extract tables (runbooks often have procedure tables)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip(" |"):
                    texts.append(row_text)

        full_text = "\n".join(texts).strip()
        if full_text:
            documents = [
                Document(
                    content=full_text,
                    metadata={
                        "source_file": filepath.name,
                        "source_path": str(filepath),
                        "folder": filepath.parent.name,
                        "doc_type": _classify_doc_type(filepath),
                        "job_id": _extract_job_id(filepath.stem),
                        "page_number": 1,
                        "file_type": "docx",
                    },
                )
            ]
            logger.info(f"Loaded DOCX: {filepath.name}")
            return documents

        return []

    except Exception as e:
        logger.error(f"Failed to load DOCX {filepath.name}: {e}")
        return []


def load_doc(filepath: Path) -> list[Document]:
    """Load a legacy .doc file (pre-2007 Word format).

    WHY separate handler: .doc files use a binary format that python-docx
    can't read. We use antiword or textract as a fallback. Your Docs folder
    has at least one .doc file (Runbook+Search.doc).
    """
    try:
        import subprocess

        # Try antiword first (faster, cleaner output)
        result = subprocess.run(
            ["antiword", str(filepath)],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            text = result.stdout.strip()
        else:
            # Fallback: try textract
            import textract

            text = textract.process(str(filepath)).decode("utf-8").strip()

        if text:
            return [
                Document(
                    content=text,
                    metadata={
                        "source_file": filepath.name,
                        "source_path": str(filepath),
                        "folder": filepath.parent.name,
                        "doc_type": _classify_doc_type(filepath),
                        "job_id": _extract_job_id(filepath.stem),
                        "page_number": 1,
                        "file_type": "doc",
                    },
                )
            ]
        return []

    except Exception as e:
        logger.warning(
            f"Could not load .doc file {filepath.name}: {e}. "
            "Install antiword (`brew install antiword`) for .doc support."
        )
        return []


def load_txt(filepath: Path) -> list[Document]:
    """Load a plain text file."""
    try:
        text = filepath.read_text(encoding="utf-8", errors="replace").strip()
        if text:
            return [
                Document(
                    content=text,
                    metadata={
                        "source_file": filepath.name,
                        "source_path": str(filepath),
                        "folder": filepath.parent.name,
                        "doc_type": _classify_doc_type(filepath),
                        "job_id": _extract_job_id(filepath.stem),
                        "page_number": 1,
                        "file_type": "txt",
                    },
                )
            ]
        return []
    except Exception as e:
        logger.error(f"Failed to load TXT {filepath.name}: {e}")
        return []


# Map file extensions to their loader functions
LOADERS = {
    ".pdf": load_pdf,
    ".pptx": load_pptx,
    ".docx": load_docx,
    ".doc": load_doc,
    ".txt": load_txt,
}

# File extensions to skip
SKIP_EXTENSIONS = {".zip", ".ds_store", ".gitkeep"}


def load_all_documents(docs_path: str | Path) -> list[Document]:
    """Load ALL documents from the Docs directory and its subfolders.

    This is the main entry point for document ingestion. It:
    1. Recursively walks Files, Files_1, Knowledge, and Temp folders
    2. Loads each file with the appropriate loader
    3. Returns a flat list of Document objects ready for chunking

    Args:
        docs_path: Path to the Docs root directory

    Returns:
        List of Document objects with content and metadata
    """
    docs_path = Path(docs_path)
    if not docs_path.exists():
        raise FileNotFoundError(f"Docs directory not found: {docs_path}")

    all_documents = []
    skipped_files = []

    for filepath in sorted(docs_path.rglob("*")):
        if not filepath.is_file():
            continue

        ext = filepath.suffix.lower()

        # Skip non-document files
        if ext in SKIP_EXTENSIONS or filepath.name.startswith("."):
            continue

        loader = LOADERS.get(ext)
        if loader:
            docs = loader(filepath)
            all_documents.extend(docs)
        else:
            skipped_files.append(filepath.name)

    if skipped_files:
        logger.warning(
            f"Skipped {len(skipped_files)} unsupported files: {skipped_files[:5]}..."
        )

    logger.info(f"Loaded {len(all_documents)} document segments from {docs_path}")
    return all_documents
